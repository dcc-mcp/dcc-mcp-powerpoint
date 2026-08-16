"""Analyzer tests — defect deck must be flagged, healthy checks never guess."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

from dcc_mcp_powerpoint.analyze import analyze_deck

ROOT = Path(__file__).resolve().parent.parent
EXAMPLE_DECK = ROOT / "examples/output/draft-dcc-mcp-framework-intro.pptx"


def _build_defect_deck(path: Path) -> None:
    from PIL import Image

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # structure: text overflow (tiny box, long text)
    tiny = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(0.6), Inches(0.35))
    tiny.text_frame.text = "this is a long line that cannot possibly fit inside a tiny box"
    # structure: out-of-bounds (beyond the right slide edge)
    rogue = slide.shapes.add_textbox(Inches(12.5), Inches(2), Inches(2.0), Inches(0.5))
    rogue.text_frame.text = "rogue"
    # format: non-brand font
    odd = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(0.6))
    run = odd.text_frame.paragraphs[0].add_run()
    run.text = "Comic Sans text"
    run.font.name = "Comic Sans MS"
    # format: low contrast (white text on white fill)
    low = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(3), Inches(0.6))
    low.fill.solid()
    low.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    lrun = low.text_frame.paragraphs[0].add_run()
    lrun.text = "invisible text"
    lrun.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # content: picture without alt text
    png = path.parent / "pixel.png"
    Image.new("RGB", (32, 32), (77, 157, 224)).save(png)
    slide.shapes.add_picture(str(png), Inches(6), Inches(4), height=Inches(0.5))
    prs.save(str(path))


def test_analyzer_flags_all_defect_classes(tmp_path: Path) -> None:
    deck = tmp_path / "defects.pptx"
    _build_defect_deck(deck)
    report = analyze_deck(deck)
    assert report["success"], report
    assert report["count"] >= 5, report
    buckets = {issue["bucket"] for issue in report["issues"]}
    assert {"structure", "format", "content"} <= buckets, report["per_bucket"]
    for issue in report["issues"]:
        assert issue["path"].startswith("/slide["), issue
    structure = [i for i in report["issues"] if i["bucket"] == "structure"]
    overflow = [i for i in structure if "overflow" in i["message"]]
    bounds = [i for i in structure if "edge" in i["message"]]
    assert overflow and overflow[0]["hint"], "overflow issue missing with hint"
    assert bounds, "out-of-bounds issue missing"
    fmt = [i for i in report["issues"] if i["bucket"] == "format"]
    assert any("Comic Sans" in i["message"] for i in fmt), "brand-font issue missing"
    assert any("contrast" in i["message"] for i in fmt), "contrast issue missing"
    content = [i for i in report["issues"] if i["bucket"] == "content"]
    assert any("alt text" in i["message"] for i in content), "alt-text issue missing"


def test_analyzer_handles_missing_file() -> None:
    report = analyze_deck(tmp_path_fixture := Path("F:/github/dcc-mcp-PowerPoint") / "does-not-exist.pptx")
    assert not report["success"]
    assert "not found" in report["reason"]
    _ = tmp_path_fixture


def test_analyzer_runs_on_the_showcase_deck() -> None:
    if not EXAMPLE_DECK.is_file():
        import pytest

        pytest.skip("showcase deck not generated")
    report = analyze_deck(EXAMPLE_DECK)
    assert report["success"], report
    # The showcase deck is validated clean structurally; analysis must not
    # crash and every issue must carry a path.
    for issue in report["issues"]:
        assert issue["path"].startswith("/slide["), issue
