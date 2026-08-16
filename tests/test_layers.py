"""Smart layer system tests — Office-free, deterministic.

Uses tests/_tmphelper.py instead of pytest's tmp_path: the DSH sandbox
denies pytest's extended-path temp directories (see helper docstring).
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from dcc_mcp_powerpoint.compiler import compile_deck
from dcc_mcp_powerpoint.deck_ir import load_deck_ir
from dcc_mcp_powerpoint.layers import (
    BUILTIN_LAYERS,
    LAYER_TAG,
    assign_layer,
    list_layers,
    recolor_layer,
    reorder_layer,
    set_layer_visibility,
    shape_layer,
    shape_tree,
    split_layer_name,
    tag_shape,
)

from ._tmphelper import make_tmp_dir, remove_tmp_dir

EXAMPLES = Path(__file__).resolve().parent.parent / "examples"


@pytest.fixture()
def deck() -> Path:
    directory = make_tmp_dir("layers-deck")
    envelope = load_deck_ir(EXAMPLES / "dcc_mcp_framework_intro.json")
    yield compile_deck(envelope, directory / "deck.pptx")
    remove_tmp_dir(directory)


@pytest.fixture()
def prs(deck: Path):
    return Presentation(str(deck))


def test_split_layer_name_roundtrip() -> None:
    assert split_layer_name("TextBox 4 ::layer=content") == ("TextBox 4", "content")
    assert split_layer_name("no tag") == ("no tag", None)
    assert split_layer_name("trailing ::layer=") == ("trailing", "")


def test_tag_shape_idempotent(prs) -> None:
    slide = prs.slides[0]
    shape = slide.shapes[0]
    assert tag_shape(shape, "media").endswith(f"{LAYER_TAG}media")
    assert shape_layer(shape) == "media"
    # Re-tagging replaces the layer, never stacks tags.
    tag_shape(shape, "content")
    assert shape_layer(shape) == "content"
    assert shape.name.count(LAYER_TAG) == 1


def test_tag_rejects_invalid_layer_names(prs) -> None:
    with pytest.raises(ValueError, match="invalid layer name"):
        tag_shape(prs.slides[0].shapes[0], "Bad Name!")
    with pytest.raises(ValueError, match="invalid layer name"):
        tag_shape(prs.slides[0].shapes[0], "1lead-digit")


def test_generated_deck_is_fully_layer_tagged(prs) -> None:
    report = list_layers(prs, "deck.pptx")
    assert report["untagged"] == []
    names = {layer["name"] for layer in report["layers"]}
    assert names == set(BUILTIN_LAYERS)
    for layer in report["layers"]:
        assert layer["total"] > 0
        for per_slide in layer["slides"]:
            assert per_slide["count"] == len(per_slide["shapes"])


def test_set_layer_visibility_hides_and_restores(prs) -> None:
    report = set_layer_visibility(prs, "decoration", False)
    assert report["changed"] > 0
    hidden = [s for slide in prs.slides for s in slide.shapes if shape_layer(s) == "decoration"]
    assert hidden and all(s._element.get("hidden") == "1" for s in hidden)
    report = set_layer_visibility(prs, "decoration", True)
    assert report["changed"] > 0
    assert all(s._element.get("hidden") is None for s in hidden)


def test_set_layer_visibility_scoped_to_slides(prs) -> None:
    report = set_layer_visibility(prs, "header", False, slides=[1])
    assert set(report["per_slide"]) == {1}
    first = [s for s in prs.slides[0].shapes if shape_layer(s) == "header"]
    other = [s for s in prs.slides[1].shapes if shape_layer(s) == "header"]
    assert all(s._element.get("hidden") == "1" for s in first)
    assert all(s._element.get("hidden") is None for s in other)


def _z_layers(slide) -> list[str | None]:
    return [shape_layer(s) for s in slide.shapes]


def _slide_index(prs, slide) -> int:
    return list(prs.slides).index(slide) + 1


def test_reorder_layer_front_and_back(prs) -> None:
    slide = next(s for s in prs.slides if any(shape_layer(sh) == "decoration" for sh in s.shapes))
    index = _slide_index(prs, slide)
    assert _z_layers(slide)[0] == "background"
    reorder_layer(prs, "decoration", "front", slides=[index])
    assert _z_layers(slide)[-1] == "decoration"
    reorder_layer(prs, "decoration", "back", slides=[index])
    assert _z_layers(slide)[0] == "decoration"


def test_reorder_layer_above_and_below(prs) -> None:
    slide = next(s for s in prs.slides if any(shape_layer(sh) == "accent" for sh in s.shapes))
    index = _slide_index(prs, slide)
    reorder_layer(prs, "accent", "above header", slides=[index])
    layers = _z_layers(slide)
    accent_last = max(i for i, l in enumerate(layers) if l == "accent")
    header_last = max(i for i, l in enumerate(layers) if l == "header")
    assert accent_last > header_last
    reorder_layer(prs, "accent", "below header", slides=[index])
    layers = _z_layers(slide)
    accent_first = min(i for i, l in enumerate(layers) if l == "accent")
    header_first = min(i for i, l in enumerate(layers) if l == "header")
    assert accent_first < header_first


def test_reorder_rejects_bad_position(prs) -> None:
    with pytest.raises(ValueError, match="position must be"):
        reorder_layer(prs, "accent", "sideways")


def test_recolor_layer_retints_solid_fills(prs) -> None:
    report = recolor_layer(prs, "accent", fill="E4572E")
    assert report["applied"] > 0
    assert report["details"]["fill"] > 0
    accents = [s for slide in prs.slides for s in slide.shapes if shape_layer(s) == "accent"]
    recolored = 0
    for shape in accents:
        try:
            if str(shape.fill.fore_color.rgb) == "E4572E":
                recolored += 1
        except (AttributeError, TypeError, ValueError):
            pass
    assert recolored == report["details"]["fill"]


def test_recolor_layer_requires_one_channel(prs) -> None:
    with pytest.raises(ValueError, match="at least one"):
        recolor_layer(prs, "accent")
    with pytest.raises(ValueError, match="invalid hex color"):
        recolor_layer(prs, "accent", fill="red")


def _picture_shapes(prs):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield index, shape


def test_assign_layer_by_selector(prs) -> None:
    report = assign_layer(prs, "media", shape="Picture")
    assert report["tagged"] > 0
    assert all(shape_layer(s) == "media" for _, s in _picture_shapes(prs))


def test_assign_layer_requires_selector(prs) -> None:
    with pytest.raises(ValueError, match="at least one selector"):
        assign_layer(prs, "media")


def test_assign_layer_no_match_raises(prs) -> None:
    with pytest.raises(ValueError, match="no shapes matched"):
        assign_layer(prs, "media", shape="Definitely Not A Shape Name")


def test_shape_tree_reports_addressing_info(prs) -> None:
    report = shape_tree(prs, "deck.pptx", slide=2)
    assert len(report["slides"]) == 1
    entry = report["slides"][0]
    assert entry["slide"] == 2
    for shape in entry["shapes"]:
        assert {"index", "name", "id", "layer", "type", "text", "hidden"} <= set(shape)
