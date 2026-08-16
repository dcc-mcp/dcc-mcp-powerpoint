"""Open XML compiler tests — Office-free, deterministic geometry."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from dcc_mcp_powerpoint.compiler import LAYOUTS, compile_deck
from dcc_mcp_powerpoint.deck_ir import (
    DeckEnvelope,
    Metadata,
    PresentationIr,
    Slide,
    load_deck_ir,
)

EXAMPLES = Path(__file__).resolve().parent.parent / "examples"


def test_layout_registry_covers_example_deck() -> None:
    envelope = load_deck_ir(EXAMPLES / "dcc_mcp_framework_intro.json")
    used = {s.semantic_layout for s in envelope.document.slides}
    missing = used - set(LAYOUTS)
    assert not missing, f"layouts missing from registry: {missing}"


def test_compile_example_deck(tmp_path: Path) -> None:
    envelope = load_deck_ir(EXAMPLES / "dcc_mcp_framework_intro.json")
    out = compile_deck(envelope, tmp_path / "deck.pptx")
    assert out.is_file() and out.stat().st_size > 10_000
    reopened = Presentation(str(out))
    assert len(reopened.slides) == len(envelope.document.slides)
    # Speaker notes survive the compile step.
    first = reopened.slides[0]
    assert first.has_notes_slide
    assert "开场" in first.notes_slide.notes_text_frame.text
    # Title cover carries the packaged master logo (brand://dcc-mcp/*).
    pictures = [s for s in first.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1


def test_compile_rejects_unknown_layout(tmp_path: Path) -> None:
    envelope = DeckEnvelope(
        schema_version="office-ir/1.0",
        kind="presentation",
        document_id="draft:bad",
        metadata=Metadata(title="bad"),
        document=PresentationIr(slides=(Slide(semantic_layout="does_not_exist"),)),
    )
    with pytest.raises(KeyError, match="does_not_exist"):
        compile_deck(envelope, tmp_path / "bad.pptx")


def test_all_example_slides_compile(tmp_path: Path) -> None:
    """Every layout in the registry compiles to a slide without error."""
    envelope = DeckEnvelope(
        schema_version="office-ir/1.0",
        kind="presentation",
        document_id="draft:layouts",
        metadata=Metadata(title="layouts"),
        document=PresentationIr(
            slides=tuple(
                Slide(
                    semantic_layout=layout,
                    title=f"L {layout}",
                    content_blocks=({"type": "bullets", "items": [f"{layout} item {i}" for i in range(3)]},),
                )
                for layout in sorted(LAYOUTS)
            )
        ),
    )
    out = compile_deck(envelope, tmp_path / "layouts.pptx")
    reopened = Presentation(str(out))
    assert len(reopened.slides) == len(LAYOUTS)


def test_image_left_text_right_places_real_image(tmp_path: Path) -> None:
    """The image layout embeds the referenced picture (real-file path)."""
    from PIL import Image

    image = tmp_path / "shot.png"
    Image.new("RGB", (320, 180), (77, 157, 224)).save(image)
    envelope = DeckEnvelope(
        schema_version="office-ir/1.0",
        kind="presentation",
        document_id="draft:image",
        metadata=Metadata(title="image"),
        document=PresentationIr(
            slides=(
                Slide(
                    semantic_layout="image_left_text_right",
                    title="Shot 010",
                    content_blocks=(
                        {"type": "image", "resource": str(image)},
                        {"type": "bullets", "items": ["版本: v12", "制作: li.ming"]},
                    ),
                ),
            )
        ),
    )
    out = compile_deck(envelope, tmp_path / "image.pptx")
    slide = Presentation(str(out)).slides[0]
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1


def test_image_layout_skips_missing_resource(tmp_path: Path) -> None:
    """A missing image degrades to a bullets note, never a broken picture."""
    envelope = DeckEnvelope(
        schema_version="office-ir/1.0",
        kind="presentation",
        document_id="draft:image",
        metadata=Metadata(title="image"),
        document=PresentationIr(
            slides=(
                Slide(
                    semantic_layout="image_left_text_right",
                    title="Shot 010",
                    content_blocks=(
                        {"type": "image", "resource": str(tmp_path / "nope.png")},
                        {"type": "bullets", "items": ["missing_asset: 渲染产物缺失"]},
                    ),
                ),
            )
        ),
    )
    out = compile_deck(envelope, tmp_path / "image.pptx")
    slide = Presentation(str(out)).slides[0]
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pictures == []
