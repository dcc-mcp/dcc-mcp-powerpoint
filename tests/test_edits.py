"""Deck patch engine tests — ppt-patch/1.0 contract, all-or-nothing saves.

Uses tests/_tmphelper.py instead of pytest's tmp_path (DSH sandbox).
"""

from __future__ import annotations

import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from dcc_mcp_powerpoint.compiler import compile_deck
from dcc_mcp_powerpoint.deck_ir import load_deck_ir
from dcc_mcp_powerpoint.edits import (
    PATCH_VERSION,
    PatchApplyError,
    PatchValidationError,
    apply_patch,
    parse_patch,
)

from ._tmphelper import make_tmp_dir, remove_tmp_dir

EXAMPLES = Path(__file__).resolve().parent.parent / "examples"


@pytest.fixture()
def deck() -> Path:
    directory = make_tmp_dir("edits-deck")
    envelope = load_deck_ir(EXAMPLES / "dcc_mcp_framework_intro.json")
    yield compile_deck(envelope, directory / "deck.pptx")
    remove_tmp_dir(directory)


def _patch(operations: list[dict]) -> dict:
    return {"schema_version": PATCH_VERSION, "metadata": {"title": "Test"}, "operations": operations}


def test_parse_patch_accepts_valid() -> None:
    parsed = parse_patch(_patch([{"op": "set_text", "select": {"layer": "header"}, "text": "x"}]))
    assert parsed["operations"][0]["op"] == "set_text"


@pytest.mark.parametrize(
    ("raw", "match"),
    [
        ({"operations": []}, "schema_version"),
        ({"schema_version": "ppt-patch/9.9", "operations": []}, "ppt-patch/1.0"),
        ({"schema_version": PATCH_VERSION}, "operations"),
        ({"schema_version": PATCH_VERSION, "operations": [{"op": "explode"}]}, "unknown operation"),
        ({"schema_version": PATCH_VERSION, "operations": [{"op": "set_text", "select": {}}]}, "requires 'text'"),
    ],
)
def test_parse_patch_rejects_invalid(raw: dict, match: str) -> None:
    with pytest.raises(PatchValidationError, match=match):
        parse_patch(raw)


def test_set_text_by_layer(deck: Path) -> None:
    out = deck.with_name("v2.pptx")
    report = apply_patch(deck, _patch([{"op": "set_text", "select": {"slide": 2, "layer": "header"}, "text": "EDITED"}]), out)
    assert report["success"] and report["operations"][0]["affected"] == 1
    prs = Presentation(str(out))
    headers = [s for s in prs.slides[1].shapes if "::layer=header" in s.name]
    assert headers and headers[0].text_frame.text == "EDITED"


def test_set_text_requires_text_frame(deck: Path) -> None:
    out = deck.with_name("v2.pptx")
    with pytest.raises(PatchApplyError, match="text frame"):
        apply_patch(
            deck,
            _patch([{"op": "set_text", "select": {"slide": 1, "layer": "decoration", "shape": "Picture"}, "text": "x"}]),
            out,
        )
    assert not out.exists()


def test_set_notes(deck: Path) -> None:
    out = deck.with_name("v2.pptx")
    apply_patch(deck, _patch([{"op": "set_notes", "slide": 2, "text": "notes v2"}]), out)
    prs = Presentation(str(out))
    assert prs.slides[1].notes_slide.notes_text_frame.text == "notes v2"


def test_set_shape_visible_all_slides(deck: Path) -> None:
    out = deck.with_name("v2.pptx")
    report = apply_patch(deck, _patch([{"op": "set_shape_visible", "select": {"slide": "all", "layer": "decoration"}, "visible": False}]), out)
    assert report["operations"][0]["affected"] > 0
    prs = Presentation(str(out))
    hidden = [s for slide in prs.slides for s in slide.shapes if "::layer=decoration" in s.name]
    assert hidden and all(s._element.get("hidden") == "1" for s in hidden)


def test_add_delete_move_slides(deck: Path) -> None:
    out = deck.with_name("v2.pptx")
    report = apply_patch(
        deck,
        _patch(
            [
                {"op": "add_slide", "slide_ir": {"semantic_layout": "bullets", "title": "Added", "content_blocks": [{"type": "bullets", "items": ["a", "b"]}]}},
                {"op": "move_slide", "slide": 17, "to": "start"},
                {"op": "delete_slide", "slide": 2},
            ]
        ),
        out,
    )
    assert report["slides_after"] == 16
    prs = Presentation(str(out))
    titles = [s.text_frame.text for slide in prs.slides for s in slide.shapes if "::layer=header" in s.name]
    assert "Added" in titles


def test_delete_shape(deck: Path) -> None:
    prs = Presentation(str(deck))
    before = len(list(prs.slides[2].shapes))
    out = deck.with_name("v2.pptx")
    report = apply_patch(deck, _patch([{"op": "delete_shape", "select": {"slide": 3, "layer": "decoration"}}]), out)
    assert report["operations"][0]["affected"] > 0
    prs = Presentation(str(out))
    assert len(list(prs.slides[2].shapes)) == before - report["operations"][0]["affected"]


def test_set_shape_fill(deck: Path) -> None:
    out = deck.with_name("v2.pptx")
    apply_patch(deck, _patch([{"op": "set_shape_fill", "select": {"slide": 1, "layer": "background"}, "color": "112233"}]), out)
    prs = Presentation(str(out))
    bg = next(s for s in prs.slides[0].shapes if "::layer=background" in s.name)
    assert str(bg.fill.fore_color.rgb) == "112233"


def test_set_image_replaces_picture(deck: Path) -> None:
    from PIL import Image

    directory = deck.parent
    png = directory / "tiny.png"
    Image.new("RGB", (64, 64), (200, 30, 30)).save(png)
    # Rewrite an image_grid slide to reference the local png so the deck
    # definitely embeds one picture.
    raw = json.loads((EXAMPLES / "dcc_mcp_framework_intro.json").read_text(encoding="utf-8"))
    for slide in raw["document"]["slides"]:
        if slide["semantic_layout"] == "image_grid":
            slide["images"] = [{"id": "pic", "uri": str(png)}]
    from dcc_mcp_powerpoint.deck_ir import parse_envelope

    image_deck = directory / "img.pptx"
    compile_deck(parse_envelope(raw), image_deck)

    png2 = directory / "tiny2.png"
    Image.new("RGB", (64, 64), (30, 30, 200)).save(png2)
    out = directory / "v2.pptx"
    report = apply_patch(image_deck, _patch([{"op": "set_image", "select": {"layer": "content", "shape": "Picture"}, "resource": str(png2), "alt": "new alt"}]), out)
    assert report["operations"][0]["affected"] >= 1
    prs = Presentation(str(out))
    from pptx.oxml.ns import qn

    replaced = [
        s
        for slide in prs.slides
        for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE and "::layer=content" in s.name
    ]
    assert replaced
    c_nv_pr = replaced[0]._element.nvPicPr.find(qn("p:cNvPr"))
    assert c_nv_pr.get("descr") == "new alt"


def test_assign_layer_operation(deck: Path) -> None:
    out = deck.with_name("v2.pptx")
    report = apply_patch(deck, _patch([{"op": "assign_layer", "layer": "media", "select": {"shape": "Picture"}}]), out)
    assert report["operations"][0]["affected"] > 0
    prs = Presentation(str(out))
    assert all("::layer=media" in s.name for slide in prs.slides for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)


def test_failed_patch_leaves_input_untouched(deck: Path) -> None:
    original = deck.read_bytes()
    with pytest.raises(PatchApplyError, match="out of range"):
        apply_patch(deck, _patch([{"op": "delete_slide", "slide": 99}]))
    assert deck.read_bytes() == original


def test_default_output_name(deck: Path) -> None:
    report = apply_patch(deck, _patch([{"op": "set_notes", "slide": 1, "text": "x"}]))
    assert report["output"] == str(deck.with_name("deck-patched.pptx"))
    assert Path(report["output"]).is_file()


def test_select_requires_key(deck: Path) -> None:
    with pytest.raises(PatchApplyError, match="at least one"):
        apply_patch(deck, _patch([{"op": "set_text", "select": {}, "text": "x"}]), deck.with_name("v2.pptx"))
