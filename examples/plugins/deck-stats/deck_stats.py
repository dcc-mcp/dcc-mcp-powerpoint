"""deck-stats — demo PowerPoint plugin (plugin.json contract).

Reads {"context": {...}, "params": {...}} on stdin and prints one JSON
result on stdout. Uses only the dcc_mcp_powerpoint public API.
"""

from __future__ import annotations

import json
import sys
from pathlib import Path


def main() -> int:
    raw = sys.stdin.read()
    try:
        request = json.loads(raw) if raw.strip() else {}
    except json.JSONDecodeError as exc:
        print(json.dumps({"success": False, "message": f"stdin is not JSON: {exc}", "context": {}}, ensure_ascii=False))
        return 1
    context = request.get("context", {})
    params = request.get("params", {})
    pptx = params.get("pptx") or context.get("pptx")
    if not pptx or not Path(pptx).is_file():
        print(json.dumps({"success": False, "message": f"pptx not found: {pptx}", "context": {}}, ensure_ascii=False))
        return 1

    from pptx import Presentation

    from dcc_mcp_powerpoint.layers import list_layers, shape_layer

    prs = Presentation(pptx)
    per_slide = []
    for index, slide in enumerate(prs.slides, start=1):
        layers: dict[str, int] = {}
        for shape in slide.shapes:
            layer = shape_layer(shape) or "untagged"
            layers[layer] = layers.get(layer, 0) + 1
        per_slide.append({"slide": index, "shapes": len(slide.shapes), "layers": layers})
    report = list_layers(prs, str(pptx))
    summary = {
        "deck": pptx,
        "slides": len(per_slide),
        "layers": {entry["name"]: entry["total"] for entry in report["layers"]},
        "per_slide": per_slide,
    }
    print(json.dumps({"success": True, "message": f"{len(per_slide)} slides inventoried", "context": summary}, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
