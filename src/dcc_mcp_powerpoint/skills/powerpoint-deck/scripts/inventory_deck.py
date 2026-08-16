"""powerpoint-deck / inventory_deck — inspect a PPTX before editing or repair.

Read-only inventory (document-pptx learnings: inspect templates before
automation; generated decks must stay inspectable/repairable): slide count,
semantic layouts used, notes coverage, picture count with alt-text gaps,
and the slide size.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


def _force_utf8_stdio() -> None:
    """Deterministic output contract: stdout/stderr are always UTF-8."""
    for stream in (sys.stdout, sys.stderr):
        if hasattr(stream, "reconfigure"):
            stream.reconfigure(encoding="utf-8", errors="replace")


def run(params: dict) -> None:
    path = Path(params["input"])
    if not path.is_file():
        print(json.dumps({"success": False, "message": f"input not found: {path}", "context": {}}, ensure_ascii=False))
        return
    prs = Presentation(str(path))
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        alt_gaps = 0
        for pic in pictures:
            c_nv_pr = pic._element.nvPicPr.find(qn("p:cNvPr"))
            if c_nv_pr is None or not c_nv_pr.get("descr"):
                alt_gaps += 1
        slides.append(
            {
                "index": index,
                "shapes": len(slide.shapes),
                "pictures": len(pictures),
                "pictures_without_alt": alt_gaps,
                "has_notes": slide.has_notes_slide and bool(slide.notes_slide.notes_text_frame.text.strip()),
            }
        )
    context = {
        "slides": len(slides),
        "size": {"width": prs.slide_width, "height": prs.slide_height},
        "notes_coverage": f"{sum(1 for s in slides if s['has_notes'])}/{len(slides)}",
        "alt_gaps": sum(s["pictures_without_alt"] for s in slides),
        "per_slide": slides,
    }
    print(json.dumps({"success": True, "message": f"inventoried {len(slides)} slides", "context": context}, ensure_ascii=False))


def main() -> None:
    _force_utf8_stdio()
    params: dict = {}
    if not sys.stdin.isatty():
        raw = sys.stdin.read()
        if raw.strip():
            try:
                params = json.loads(raw)
            except json.JSONDecodeError:
                params = {}
    if not params:
        parser = argparse.ArgumentParser(description="Inventory a PPTX (read-only)")
        parser.add_argument("--input", required=True)
        params = vars(parser.parse_args())
    try:
        run(params)
    except Exception as exc:  # noqa: BLE001
        print(json.dumps({"success": False, "message": str(exc), "context": {}}, ensure_ascii=False))
        sys.exit(1)


if __name__ == "__main__":
    main()
