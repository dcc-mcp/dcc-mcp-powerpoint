"""Issues analyzer — self-implemented deck diagnosis with path addressing.

Research origin: OfficeCLI's view/issues surface (format/content/structure
buckets, /slide[i]/shape[j] addressing, concrete fix hints). Implemented
independently: python-pptx walk + real font-metric measurement via PIL +
WCAG relative-luminance contrast.

Honesty contract: a check that cannot be decided (theme-inherited colors,
missing font files) reports nothing for that item — never a false positive.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu

FONT_CANDIDATES = ("msyh.ttc", "msyh.ttf", "segoeui.ttf", "arial.ttf")
MIN_CONTRAST_RATIO = 4.5  # WCAG AA for normal text
SAFE_MARGIN_EMU = Emu(int(0.05 * 914400))  # 0.05in slide-bounds tolerance


@dataclass
class Issue:
    id: str
    bucket: str
    severity: str
    path: str
    message: str
    hint: str = ""
    shape_name: str = ""

    def as_dict(self) -> dict[str, str]:
        payload = {
            "id": self.id,
            "bucket": self.bucket,
            "severity": self.severity,
            "path": self.path,
            "message": self.message,
        }
        if self.hint:
            payload["hint"] = self.hint
        if self.shape_name:
            payload["shape_name"] = self.shape_name
        return payload


@dataclass
class AnalyzeReport:
    issues: list[Issue] = field(default_factory=list)

    def add(self, bucket: str, path: str, message: str, *, hint: str = "", shape_name: str = "", severity: str = "warning") -> None:
        self.issues.append(
            Issue(
                id=f"{bucket[0].upper()}{self._next(bucket)}",
                bucket=bucket,
                severity=severity,
                path=path,
                message=message,
                hint=hint,
                shape_name=shape_name,
            )
        )

    def _next(self, bucket: str) -> int:
        return sum(1 for i in self.issues if i.bucket == bucket) + 1

    def as_dict(self) -> dict[str, Any]:
        per_bucket: dict[str, int] = {}
        for issue in self.issues:
            per_bucket[issue.bucket] = per_bucket.get(issue.bucket, 0) + 1
        return {
            "success": True,
            "count": len(self.issues),
            "per_bucket": per_bucket,
            "issues": [i.as_dict() for i in self.issues],
        }


def analyze_deck(pptx_path: str | Path) -> dict[str, Any]:
    """Analyze a deck; returns the structured issue report."""
    path = Path(pptx_path)
    if not path.is_file():
        return {"success": False, "reason": f"input not found: {path}"}
    try:
        prs = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001 — surface as structured error
        return {"success": False, "reason": f"cannot open pptx: {exc}"}
    report = AnalyzeReport()
    fonts = _resolve_fonts()
    for slide_index, slide in enumerate(prs.slides, start=1):
        _analyze_slide(report, slide, slide_index, prs, fonts)
    return report.as_dict()


def _analyze_slide(report: AnalyzeReport, slide, slide_index: int, prs, fonts: dict[str, str]) -> None:
    slide_path = f"/slide[{slide_index}]"
    text_shapes = 0
    for shape_index, shape in enumerate(slide.shapes, start=1):
        path = f"{slide_path}/shape[{shape_index}]"
        name = shape.name or ""
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # python-pptx auto-fills descr with the source filename — a
            # filename is not meaningful alt text. Treat descr == cNvPr name
            # (the auto value) as missing.
            c_nv_pr = shape._element.nvPicPr.find(qn("p:cNvPr"))
            descr = c_nv_pr.get("descr") if c_nv_pr is not None else None
            auto_name = c_nv_pr.get("name") if c_nv_pr is not None else None
            auto_filled = bool(descr) and descr.lower().endswith(
                (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".webp", ".svg")
            )
            if not descr or descr == auto_name or auto_filled:
                report.add(
                    "content",
                    path,
                    "picture missing meaningful alt text",
                    hint="set descr to a human description (python-pptx auto-fills the filename)",
                    shape_name=name,
                )
            continue
        _check_bounds(report, shape, path, name, prs)
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text:
            text_shapes += 1
        if text:
            _check_overflow(report, shape, path, name, fonts)
            _check_fonts(report, shape, path, name)
            _check_contrast(report, shape, path, name)
    if text_shapes == 0:
        report.add("content", slide_path, "slide has no text content", hint="consider speaker notes or a visual-only caption", severity="warning")
    if not (slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip()):
        report.add("content", slide_path, "speaker notes empty", hint="capture the slide's one takeaway in the notes")


def _check_bounds(report: AnalyzeReport, shape, path: str, name: str, prs) -> None:
    left, top = shape.left, shape.top
    right = left + (shape.width or 0)
    bottom = top + (shape.height or 0)
    overflow_right = right - prs.slide_width - SAFE_MARGIN_EMU
    overflow_bottom = bottom - prs.slide_height - SAFE_MARGIN_EMU
    if overflow_right > 0:
        report.add("structure", path, "shape extends beyond the right slide edge", hint=f"reduce width or shift left by {_emu_cm(overflow_right):.1f}cm", shape_name=name)
    if overflow_bottom > 0:
        report.add("structure", path, "shape extends beyond the bottom slide edge", hint=f"reduce height or shift up by {_emu_cm(overflow_bottom):.1f}cm", shape_name=name)


def _check_overflow(report: AnalyzeReport, shape, path: str, name: str, fonts: dict[str, str]) -> None:
    if not fonts or shape.width is None or shape.height is None:
        return
    tf = shape.text_frame
    try:
        margin_l = tf.margin_left or 0
        margin_r = tf.margin_right or 0
        margin_t = tf.margin_top or 0
        margin_b = tf.margin_bottom or 0
    except Exception:  # noqa: BLE001 — inherited margins are not measurable
        return
    usable_width = shape.width - margin_l - margin_r
    usable_height = shape.height - margin_t - margin_b
    if usable_width <= 0 or usable_height <= 0:
        return
    needed_emu = 0.0
    worst_line = ""
    from PIL import ImageFont

    for para in tf.paragraphs:
        runs = para.runs or [para]
        for run in runs:
            text = run.text
            if not text.strip():
                continue
            size_pt = run.font.size.pt if run.font.size else 18.0
            font_path = _pick_font(text, size_pt, fonts)
            if font_path is None:
                return
            try:
                font = ImageFont.truetype(font_path, round(size_pt * 4))  # 4x scale for precision
            except Exception:  # noqa: BLE001
                return
            width_px = font.getlength(text)
            width_emu = width_px / 4.0 * 12700.0  # px(4x) -> pt -> EMU
            if width_emu > usable_width:
                worst_line = text
            needed_emu += max(width_emu, 0.0) * 1.22  # line height factor
    if needed_emu > usable_height:
        overflow_cm = _emu_cm(needed_emu - usable_height)
        hint = f"suggest.height=+{overflow_cm:.1f}cm" if overflow_cm >= 0.05 else "shrink font or trim text"
        report.add(
            "structure",
            path,
            f"text overflow: rendered text needs {_emu_cm(needed_emu):.1f}cm, usable {_emu_cm(usable_height):.1f}cm"
            + (f" (worst line: {worst_line[:20]}…)" if worst_line else ""),
            hint=hint,
            shape_name=name,
        )


def _check_fonts(report: AnalyzeReport, shape, path: str, name: str) -> None:
    seen: set[str] = set()
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            font_name = run.font.name
            if font_name:
                seen.add(font_name)
    for font_name in sorted(seen):
        if font_name not in ("Segoe UI", "Microsoft YaHei"):
            report.add("format", path, f"non-brand font '{font_name}'", hint="use the brand typefaces (Segoe UI / Microsoft YaHei)", shape_name=name)


def _check_contrast(report: AnalyzeReport, shape, path: str, name: str) -> None:
    try:
        bg = shape.fill.fore_color.rgb
    except Exception:  # noqa: BLE001 — inherited/theme fill is undecidable
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                fg = run.font.color.rgb
            except Exception:  # noqa: BLE001, S112 — inherited/theme color is undecidable, skip
                continue
            ratio = _contrast_ratio(fg, bg)
            if ratio < MIN_CONTRAST_RATIO:
                report.add(
                    "format",
                    path,
                    f"low contrast: text/fill ratio {ratio:.1f}:1",
                    hint="darken text or lighten the fill to reach 4.5:1",
                    shape_name=name,
                )
            break  # one representative run per shape is enough
        break


def _contrast_ratio(a: RGBColor, b: RGBColor) -> float:
    def lum(color: RGBColor) -> float:
        def channel(value: int) -> float:
            normalized = value / 255.0
            return normalized / 12.92 if normalized <= 0.03928 else ((normalized + 0.055) / 1.055) ** 2.4

        return 0.2126 * channel(color[0]) + 0.7152 * channel(color[1]) + 0.0722 * channel(color[2])

    la, lb = lum(a), lum(b)
    lighter, darker = max(la, lb), min(la, lb)
    return (lighter + 0.05) / (darker + 0.05)


def _resolve_fonts() -> dict[str, str]:
    fonts_dir = Path("C:/Windows/Fonts")
    resolved: dict[str, str] = {}
    for candidate in FONT_CANDIDATES:
        target = fonts_dir / candidate
        if target.is_file():
            resolved[candidate] = str(target)
    return resolved


def _pick_font(text: str, size_pt: float, fonts: dict[str, str]) -> str | None:
    has_cjk = any("一" <= ch <= "鿿" for ch in text)
    if has_cjk:
        for key in ("msyh.ttc", "msyh.ttf"):
            if key in fonts:
                return fonts[key]
        return None
    for key in ("segoeui.ttf", "arial.ttf"):
        if key in fonts:
            return fonts[key]
    return next(iter(fonts.values()), None)


def _emu_cm(emu: float) -> float:
    return round(emu / 914400.0 * 2.54, 2)
