"""Open XML compiler — Deck IR → PPTX (16:9).

Contract-first: the compiler is an *implementation* of the deck contract.
Semantic layouts are a registry (open for extension — adding a layout means
adding one entry, never editing dispatch logic). It renders structure only:
final layout fidelity, PDF and previews come from the COM renderer.

Design system (showcase quality):
- unified margins + title header with accent rule on every content slide
- editorial bullet rows (accent dot + hairline) instead of box-heavy slides
- kpi cards with accent top strip, architecture with numbered layer chips,
  timeline with alternating above/below labels
- centered hero cover, ghost-number section covers, deck-title footer
"""

from __future__ import annotations

import re
from collections.abc import Callable
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from .brand import resolve_logo
from .deck_ir import DeckEnvelope, Slide

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Brand palette (default theme until brand:// templates land in the registry).
COLOR_BG = RGBColor(0x0F, 0x14, 0x1E)
COLOR_BG_SOFT = RGBColor(0x17, 0x1F, 0x2E)
COLOR_PANEL = RGBColor(0x1E, 0x2A, 0x3D)
COLOR_ACCENT = RGBColor(0x4D, 0x9D, 0xE0)
COLOR_ACCENT_2 = RGBColor(0x7F, 0xC8, 0xA9)
COLOR_TEXT = RGBColor(0xE8, 0xEC, 0xF2)
COLOR_MUTED = RGBColor(0x9A, 0xA7, 0xBC)
COLOR_GHOST = RGBColor(0x22, 0x30, 0x47)

FONT_LATIN = "Segoe UI"
FONT_CJK = "Microsoft YaHei"

# Grid
MARGIN = 0.9
CONTENT_WIDTH = 11.5
HEADER_TITLE_Y = 0.55
HEADER_RULE_Y = 1.42
BODY_TOP = 1.78


class DeckCompiler:
    """Compiles a DeckEnvelope into a PPTX file via python-pptx (Open XML)."""

    def __init__(self, envelope: DeckEnvelope) -> None:
        self.envelope = envelope
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self.blank = self.prs.slide_layouts[6]

    # -- primitives ---------------------------------------------------------

    def _set_run_font(self, run, size_pt: float, color: RGBColor, bold: bool = False) -> None:
        """Latin and East-Asian typefaces are set explicitly: Segoe UI for
        latin glyphs, Microsoft YaHei for CJK — no per-character font
        fallback in PowerPoint."""
        font = run.font
        font.name = FONT_LATIN
        font.size = Pt(size_pt)
        font.bold = bold
        font.color.rgb = color
        rpr = run._r.get_or_add_rPr()
        # Per-script font slots (OfficeCLI lesson): latin, east-asian and
        # complex-script typefaces are declared separately so CJK/RTL text
        # never falls back per-character.
        for tag, typeface in (("a:latin", FONT_LATIN), ("a:ea", FONT_CJK), ("a:cs", FONT_LATIN)):
            existing = rpr.find(qn(tag))
            if existing is None:
                existing = rpr.makeelement(qn(tag), {})
                rpr.append(existing)
            existing.set("typeface", typeface)

    def _add_textbox(self, slide, left, top, width, height, text: str, *, size: float = 18, color: RGBColor = COLOR_TEXT, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
        box = slide.shapes.add_textbox(left, top, width, height)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(4)
        lines = text.split("\n")
        for i, line in enumerate(lines):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = align
            run = para.add_run()
            run.text = line
            self._set_run_font(run, size, color, bold)
        return box

    def _add_shape(self, slide, shape_type, left, top, width, height, *, fill: RGBColor = COLOR_PANEL, line: RGBColor | None = None, radius: float | None = None):
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.shadow.inherit = False
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = line
        if radius is not None and shape.adjustments:
            shape.adjustments[0] = radius
        shape.text_frame.word_wrap = True
        return shape

    def _fill_shape(self, shape, text: str, *, size: float = 16, color: RGBColor = COLOR_TEXT, bold: bool = False, align=PP_ALIGN.CENTER) -> None:
        tf = shape.text_frame
        tf.margin_left = tf.margin_right = Pt(8)
        tf.margin_top = tf.margin_bottom = Pt(6)
        for i, line in enumerate(text.split("\n")):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = align
            run = para.add_run()
            run.text = line
            self._set_run_font(run, size, color, bold)

    def _background(self, slide, color: RGBColor = COLOR_BG) -> None:
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        rect.shadow.inherit = False
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()

    def _ghost_circle(self, slide, left, top, width, height) -> None:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height))
        circle.shadow.inherit = False
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_GHOST
        circle.line.fill.background()

    def _add_picture(self, slide, path, left, top, *, height=None, width=None, alt: str = ""):
        """Embed a picture with alt text — accessibility is part of
        authoring, not an afterthought (document-pptx learnings)."""
        pic = slide.shapes.add_picture(str(path), left, top, height=height, width=width)
        pic.shadow.inherit = False
        c_nv_pr = pic._element.nvPicPr.find(qn("p:cNvPr"))
        if c_nv_pr is not None:
            c_nv_pr.set("descr", alt)
        return pic

    def _content_header(self, slide, title: str) -> None:
        self._add_textbox(slide, Inches(MARGIN), Inches(HEADER_TITLE_Y), Inches(CONTENT_WIDTH), Inches(0.75), title, size=28, bold=True)
        self._add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(HEADER_RULE_Y), Inches(0.9), Pt(3.5)).fill.fore_color.rgb = COLOR_ACCENT

    def _footer(self, slide) -> None:
        self._add_textbox(slide, Inches(MARGIN), SLIDE_HEIGHT - Inches(0.42), Inches(8.0), Inches(0.3), self.envelope.metadata.title, size=10, color=COLOR_MUTED)

    def _notes(self, slide, notes: str | None) -> None:
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _page_number(self, slide, number: int, total: int) -> None:
        box = slide.shapes.add_textbox(SLIDE_WIDTH - Inches(1.6), SLIDE_HEIGHT - Inches(0.5), Inches(1.2), Inches(0.35))
        box.shadow.inherit = False
        box.text_frame.word_wrap = False
        para = box.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.RIGHT
        run = para.add_run()
        run.text = f"{number:02d} / {total:02d}"
        self._set_run_font(run, 10, COLOR_MUTED)

    # -- slide dispatch (OCP: one registry, one entry per layout) -----------

    def compile(self, out_path: str | Path) -> Path:
        document = self.envelope.document
        for index, slide_ir in enumerate(document.slides):
            handler = LAYOUTS.get(slide_ir.semantic_layout)
            if handler is None:
                raise KeyError(
                    f"unknown semantic_layout '{slide_ir.semantic_layout}' (slide {index + 1}); "
                    f"known: {', '.join(sorted(LAYOUTS))}"
                )
            slide = self.prs.slides.add_slide(self.blank)
            self._background(slide)
            handler(self, slide, slide_ir)
            self._footer(slide)
            self._notes(slide, slide_ir.speaker_notes)
            self._page_number(slide, index + 1, len(document.slides))
        out = Path(out_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(out))
        return out


# -- layouts ---------------------------------------------------------------

def _hero_logo(c: DeckCompiler, slide) -> None:
    logo = resolve_logo(c.envelope.template.uri if c.envelope.template else None, dark_background=True)
    if logo is not None:
        c._add_picture(slide, logo, Inches(5.92), Inches(0.75), height=Inches(0.6), alt="DCC-MCP logo")


def _title_cover(c: DeckCompiler, slide, ir: Slide) -> None:
    c._background(slide, COLOR_BG)
    c._ghost_circle(slide, 10.15, 4.85, 3.05, 2.5)
    _hero_logo(c, slide)
    c._add_textbox(slide, Inches(1.5), Inches(2.9), Inches(10.3), Inches(1.4), ir.title or c.envelope.metadata.title, size=54, bold=True, align=PP_ALIGN.CENTER)
    c._add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.22), Inches(4.35), Inches(0.9), Pt(4)).fill.fore_color.rgb = COLOR_ACCENT
    subtitle = next((b for b in ir.content_blocks if b["type"] == "text"), None)
    if subtitle:
        text = "\n".join(subtitle.get("paragraphs", []))
        c._add_textbox(slide, Inches(1.5), Inches(4.65), Inches(10.3), Inches(1.0), text, size=18, color=COLOR_MUTED, align=PP_ALIGN.CENTER)
    c._add_textbox(slide, Inches(1.5), Inches(6.85), Inches(10.3), Inches(0.35), f"{c.envelope.metadata.title} · {c.envelope.document_id}", size=11, color=COLOR_MUTED, align=PP_ALIGN.CENTER)


def _section_cover(c: DeckCompiler, slide, ir: Slide) -> None:
    c._background(slide, COLOR_BG_SOFT)
    c._ghost_circle(slide, 10.5, 0.35, 2.7, 2.7)
    title = ir.title or ""
    number = ""
    match = re.match(r"(\d+)\s*[·.:]\s*(.*)", title)
    if match:
        number, title = match.group(1), match.group(2)
    if number:
        c._add_textbox(slide, Inches(0.7), Inches(0.7), Inches(5.0), Inches(4.2), number, size=200, color=COLOR_GHOST, bold=True)
    c._add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(3.05), Inches(0.18), Inches(1.2)).fill.fore_color.rgb = COLOR_ACCENT_2
    c._add_textbox(slide, Inches(1.3), Inches(3.0), Inches(11.0), Inches(1.4), title, size=44, bold=True)
    text = next((b for b in ir.content_blocks if b["type"] == "text"), None)
    if text:
        c._add_textbox(slide, Inches(1.35), Inches(4.5), Inches(10.5), Inches(0.6), "\n".join(text.get("paragraphs", [])), size=15, color=COLOR_MUTED)


def _bullets(c: DeckCompiler, slide, ir: Slide) -> None:
    if ir.title:
        c._content_header(slide, ir.title)
    y = BODY_TOP
    for block in ir.content_blocks:
        if block["type"] != "bullets":
            continue
        items = block.get("items", [])
        for i, item in enumerate(items):
            c._add_shape(slide, MSO_SHAPE.OVAL, Inches(MARGIN + 0.05), Inches(y + 0.17), Pt(13), Pt(13), fill=COLOR_ACCENT_2)
            c._add_textbox(slide, Inches(MARGIN + 0.45), Inches(y), Inches(CONTENT_WIDTH - 0.55), Inches(0.9), item, size=16.5)
            if i < len(items) - 1:
                c._add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(y + 0.95), Inches(CONTENT_WIDTH), Pt(1.2)).fill.fore_color.rgb = COLOR_GHOST
            y += 1.02


def _two_columns(c: DeckCompiler, slide, ir: Slide) -> None:
    if ir.title:
        c._content_header(slide, ir.title)
    texts = [b for b in ir.content_blocks if b["type"] == "bullets"]
    for col, block in enumerate(texts[:2]):
        left = Inches(MARGIN + col * 6.0)
        items = block.get("items", [])
        header_text = items[0] if items else ""
        head = c._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(BODY_TOP), Inches(5.6), Inches(0.62), fill=COLOR_ACCENT if col == 0 else COLOR_ACCENT_2, radius=0.16)
        c._fill_shape(head, header_text, size=16, bold=True, align=PP_ALIGN.LEFT)
        body = c._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(BODY_TOP + 0.74), Inches(5.6), Inches(4.1), fill=COLOR_PANEL, radius=0.05)
        c._fill_shape(body, "\n".join("●  " + i for i in items[1:]), size=14.5, align=PP_ALIGN.LEFT)


def _comparison(c: DeckCompiler, slide, ir: Slide) -> None:
    if ir.title:
        c._content_header(slide, ir.title)
    texts = [b for b in ir.content_blocks if b["type"] == "bullets"]
    for col, block in enumerate(texts[:2]):
        left = Inches(MARGIN + col * 6.0)
        items = block.get("items", [])
        color = COLOR_ACCENT if col == 0 else COLOR_ACCENT_2
        head = c._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(BODY_TOP), Inches(5.6), Inches(0.62), fill=color, radius=0.16)
        c._fill_shape(head, items[0] if items else "", size=16, bold=True)
        body = c._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(BODY_TOP + 0.74), Inches(5.6), Inches(4.1), fill=COLOR_PANEL, radius=0.05)
        c._fill_shape(body, "\n".join("✓  " + i for i in items[1:]), size=14.5, align=PP_ALIGN.LEFT)


def _timeline(c: DeckCompiler, slide, ir: Slide) -> None:
    if ir.title:
        c._content_header(slide, ir.title)
    items = next((b["items"] for b in ir.content_blocks if b["type"] == "bullets"), [])
    n = max(len(items), 1)
    step = CONTENT_WIDTH / n
    label_w = min(3.6, max(step - 0.3, 1.0))
    y_line = Inches(3.55)
    c._add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(MARGIN), y_line, Inches(CONTENT_WIDTH), Pt(3)).fill.fore_color.rgb = COLOR_ACCENT
    for i, item in enumerate(items):
        x = Inches(MARGIN + step * i + step / 2)
        circle = c._add_shape(slide, MSO_SHAPE.OVAL, x - Inches(0.17), y_line - Inches(0.17), Inches(0.34), Inches(0.34), fill=COLOR_BG, line=COLOR_ACCENT_2)
        circle.line.width = Pt(1.5)
        c._fill_shape(circle, f"{i + 1}", size=12, color=COLOR_ACCENT_2, bold=True)
        label_y = Inches(2.3) if i % 2 == 0 else Inches(3.9)
        c._add_textbox(slide, x - Inches(label_w / 2), label_y, Inches(label_w), Inches(1.05), item, size=13.5, align=PP_ALIGN.CENTER)


def _kpi_dashboard(c: DeckCompiler, slide, ir: Slide) -> None:
    if ir.title:
        c._content_header(slide, ir.title)
    items = next((b["items"] for b in ir.content_blocks if b["type"] == "bullets"), [])
    cols = 3
    for i, item in enumerate(items[:9]):
        row, col = divmod(i, cols)
        left = Inches(MARGIN + col * 3.9)
        top = Inches(BODY_TOP + row * 2.1)
        card = c._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.6), Inches(1.85), fill=COLOR_PANEL, radius=0.07)
        c._add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.4), top, Inches(0.55), Pt(4)).fill.fore_color.rgb = COLOR_ACCENT_2
        parts = item.split("|", 1)
        value = parts[0].strip()
        label = parts[1].strip() if len(parts) > 1 else ""
        c._fill_shape(card, value + "\n" + label, size=26, bold=(len(parts) == 1), color=COLOR_ACCENT_2 if len(parts) > 1 else COLOR_TEXT)


def _technical_architecture(c: DeckCompiler, slide, ir: Slide) -> None:
    if ir.title:
        c._content_header(slide, ir.title)
    items = next((b["items"] for b in ir.content_blocks if b["type"] == "bullets"), [])
    top_ins = BODY_TOP + 0.05
    centers: list[float] = []
    for i, item in enumerate(items):
        top = Inches(top_ins + i * 0.95)
        chip = c._add_shape(slide, MSO_SHAPE.OVAL, Inches(MARGIN + 1.45), top + Inches(0.16), Inches(0.42), Inches(0.42), fill=COLOR_BG_SOFT, line=COLOR_ACCENT)
        chip.line.width = Pt(1.2)
        c._fill_shape(chip, f"{i + 1}", size=12, color=COLOR_ACCENT, bold=True)
        box = c._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(MARGIN + 2.15), top, Inches(CONTENT_WIDTH - 2.3), Inches(0.78), fill=COLOR_PANEL, radius=0.10)
        c._fill_shape(box, item, size=14.5)
        centers.append(top_ins + i * 0.95 + 0.39)
    x = Inches(MARGIN + 1.66)
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, Inches(centers[0] + 0.18), x, Inches(centers[-1] + 0.18))
    conn.line.color.rgb = COLOR_ACCENT
    conn.shadow.inherit = False


def _image_left_text_right(c: DeckCompiler, slide, ir: Slide) -> None:
    if ir.title:
        c._content_header(slide, ir.title)
    image = next((b for b in ir.content_blocks if b["type"] == "image"), None)
    if image and Path(image.get("resource", "")).is_file():
        c._add_picture(slide, image["resource"], Inches(MARGIN), Inches(BODY_TOP), height=Inches(4.7), alt=str(image.get("resource", "")))
    texts = [b for b in ir.content_blocks if b["type"] == "bullets"]
    if texts:
        panel = c._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.05), Inches(BODY_TOP), Inches(CONTENT_WIDTH - 6.2), Inches(4.7), fill=COLOR_PANEL, radius=0.05)
        c._fill_shape(panel, "\n".join("●  " + i for i in texts[0].get("items", [])), size=14.5, align=PP_ALIGN.LEFT)


def _image_grid(c: DeckCompiler, slide, ir: Slide) -> None:
    """Card grid over slide.images: uniform picture cards with captions.

    Slide.images is the contract (Resource: id = caption, uri = path).
    Missing files render an explicit placeholder note, never a broken
    picture (same rule as image_left_text_right).
    """
    if ir.title:
        c._content_header(slide, ir.title)
    images = list(ir.images)
    if not images:
        return
    cols = 4 if len(images) >= 8 else 2
    gap = 0.25
    card_w = (CONTENT_WIDTH - gap * (cols - 1)) / cols
    rows = (len(images) + cols - 1) // cols
    card_h = min(1.72, (5.3 - gap * (rows - 1)) / rows)
    for i, res in enumerate(images):
        row, col = divmod(i, cols)
        left = MARGIN + col * (card_w + gap)
        top = BODY_TOP + row * (card_h + gap)
        c._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(card_w), Inches(card_h), fill=COLOR_PANEL, radius=0.08)
        if Path(res.uri).is_file():
            pic_h = card_h - 0.4
            c._add_picture(slide, res.uri, Inches(left + 0.1), Inches(top + 0.06), height=Inches(pic_h), alt=res.id)
        else:
            c._add_textbox(slide, Inches(left + 0.1), Inches(top + 0.3), Inches(card_w - 0.2), Inches(0.4), f"missing_asset: {res.uri}", size=11, color=COLOR_MUTED, align=PP_ALIGN.CENTER)
        c._add_textbox(slide, Inches(left), Inches(top + card_h - 0.34), Inches(card_w), Inches(0.3), res.id, size=11, color=COLOR_MUTED, align=PP_ALIGN.CENTER)


def _closing(c: DeckCompiler, slide, ir: Slide) -> None:
    c._background(slide, COLOR_BG)
    c._ghost_circle(slide, 0.35, 0.35, 2.9, 2.9)
    logo = resolve_logo(c.envelope.template.uri if c.envelope.template else None, dark_background=True)
    if logo is not None:
        c._add_picture(slide, logo, Inches(5.92), Inches(2.15), height=Inches(0.8), alt="DCC-MCP logo")
    c._add_textbox(slide, Inches(1.5), Inches(3.35), Inches(10.3), Inches(1.2), ir.title or "Thanks", size=46, bold=True, align=PP_ALIGN.CENTER)
    text = next((b for b in ir.content_blocks if b["type"] == "text"), None)
    if text:
        c._add_textbox(slide, Inches(1.5), Inches(4.75), Inches(10.3), Inches(0.7), "\n".join(text.get("paragraphs", [])), size=16, color=COLOR_MUTED, align=PP_ALIGN.CENTER)


LAYOUTS: dict[str, Callable] = {
    "title_cover": _title_cover,
    "section_cover": _section_cover,
    "bullets": _bullets,
    "two_columns": _two_columns,
    "comparison": _comparison,
    "timeline": _timeline,
    "kpi_dashboard": _kpi_dashboard,
    "technical_architecture": _technical_architecture,
    "image_left_text_right": _image_left_text_right,
    "image_grid": _image_grid,
    "closing": _closing,
}


def compile_deck(envelope: DeckEnvelope, out_path: str | Path) -> Path:
    """Compile a DeckEnvelope to PPTX via the Open XML backend."""
    return DeckCompiler(envelope).compile(out_path)
