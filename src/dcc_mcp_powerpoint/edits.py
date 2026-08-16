"""Deck patch engine — structured, all-or-nothing edits on existing PPTX.

Patch document contract (ppt-patch/1.0):

    {
      "schema_version": "ppt-patch/1.0",
      "metadata": {"title": "optional deck title for new-slide footers"},
      "operations": [
        {"op": "set_text", "select": {...}, "text": "new text"},
        ...
      ]
    }

Selectors address shapes without coordinates. Every selector key is
optional and they combine with AND; at least one is required:

    slide  1-based slide index, a list of indexes, or "all"
    layer  semantic layer tag (see layers.py), e.g. "header", "content"
    shape  exact shape name or substring
    id     shape cNvPr id
    index  1-based position of the shape on its slide

Operations:
    set_text          replace the text of every selected shape
    set_notes         replace speaker notes of one slide (op.slide, op.text)
    set_image         replace a picture (op.resource, op.alt) at its position
    set_shape_fill    solid fill color (op.color, RRGGBB) on selected shapes
    set_shape_line    solid line color (op.color, RRGGBB) on selected shapes
    set_shape_visible show/hide selected shapes (op.visible)
    delete_shape      remove selected shapes
    add_slide         compile a slide IR (op.slide_ir) with the deck compiler
                      and insert it at op.position (1-based, default "end")
    delete_slide      remove slide op.slide (1-based)
    move_slide        move slide op.slide to op.to (1-based | "start" | "end")
    assign_layer      tag selected shapes into op.layer

Transaction semantics: a failed operation aborts the patch before anything
is saved — the input file is never left half-edited. The output file is
written only when every operation succeeded.

This module is opt-in (never imported at package import time): it uses
python-pptx, a dev/test-only dependency, same policy as compiler/analyze.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Pt

from .compiler import DeckCompiler
from .deck_ir import (
    DOCUMENT_KIND,
    IR_VERSION,
    DeckEnvelope,
    Metadata,
    PresentationIr,
    parse_slide,
)
from .layers import _shape_id, assign_layer, shape_layer

PATCH_VERSION = "ppt-patch/1.0"

_OPERATIONS = frozenset(
    {
        "set_text",
        "set_notes",
        "set_image",
        "set_shape_fill",
        "set_shape_line",
        "set_shape_visible",
        "delete_shape",
        "add_slide",
        "delete_slide",
        "move_slide",
        "assign_layer",
    }
)

_OP_REQUIRED: dict[str, tuple[str, ...]] = {
    "set_text": ("select", "text"),
    "set_notes": ("slide", "text"),
    "set_image": ("select", "resource"),
    "set_shape_fill": ("select", "color"),
    "set_shape_line": ("select", "color"),
    "set_shape_visible": ("select", "visible"),
    "delete_shape": ("select",),
    "add_slide": ("slide_ir",),
    "delete_slide": ("slide",),
    "move_slide": ("slide", "to"),
    "assign_layer": ("layer", "select"),
}


class PatchValidationError(ValueError):
    """A patch document violated the ppt-patch contract. Carries a path hint."""

    def __init__(self, path: str, message: str) -> None:
        super().__init__(f"[{path}] {message}")
        self.path = path
        self.message = message


class PatchApplyError(ValueError):
    """An operation could not be applied to the target deck."""

    def __init__(self, operation_index: int, message: str) -> None:
        super().__init__(f"operation {operation_index} failed: {message}")
        self.operation_index = operation_index
        self.message = message


def parse_patch(raw: Any) -> dict[str, Any]:
    """Validate a ppt-patch/1.0 document; returns it unchanged on success."""
    if not isinstance(raw, dict):
        raise PatchValidationError("$", "patch must be an object")
    version = raw.get("schema_version")
    if version != PATCH_VERSION:
        raise PatchValidationError("$.schema_version", f"expected '{PATCH_VERSION}', got {version!r}")
    operations = raw.get("operations")
    if not isinstance(operations, list) or not operations:
        raise PatchValidationError("$.operations", "must be a non-empty list")
    for index, op in enumerate(operations):
        path = f"$.operations[{index}]"
        if not isinstance(op, dict):
            raise PatchValidationError(path, "operation must be an object")
        name = op.get("op")
        if name not in _OPERATIONS:
            raise PatchValidationError(f"{path}.op", f"unknown operation {name!r}")
        for key in _OP_REQUIRED[name]:
            if key not in op:
                raise PatchValidationError(f"{path}.{key}", f"operation '{name}' requires '{key}'")
    return dict(raw)


def _normalize_slides(prs, value: Any) -> set[int] | None:
    if value is None or value == "all":
        return None
    total = len(prs.slides._sldIdLst)
    if isinstance(value, int):
        if not 1 <= value <= total:
            raise PatchApplyError(-1, f"slide {value} out of range (deck has {total} slides)")
        return {value}
    if isinstance(value, list) and all(isinstance(i, int) for i in value):
        invalid = [i for i in value if not 1 <= i <= total]
        if invalid:
            raise PatchApplyError(-1, f"slides {invalid} out of range (deck has {total} slides)")
        return set(value)
    raise PatchApplyError(-1, f"invalid slide selector {value!r} (int, list of ints, or 'all')")


def _select(prs, select: dict[str, Any]) -> list[tuple[Any, Any]]:
    """Resolve a selector to (slide, shape) pairs; fail loudly when empty."""
    unknown = set(select) - {"slide", "slides", "layer", "shape", "id", "index"}
    if unknown:
        raise PatchApplyError(-1, f"unknown select keys: {sorted(unknown)}")
    if not set(select) & {"slide", "slides", "layer", "shape", "id", "index"}:
        raise PatchApplyError(-1, "select needs at least one of slide/slides/layer/shape/id/index")
    slides_filter = _normalize_slides(prs, select.get("slide", select.get("slides")))
    layer = select.get("layer")
    name = select.get("shape")
    shape_id = str(select["id"]) if select.get("id") is not None else None
    position = select.get("index")
    found: list[tuple[Any, Any]] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        if slides_filter is not None and slide_index not in slides_filter:
            continue
        for shape_position, shape in enumerate(slide.shapes, start=1):
            if layer is not None and shape_layer(shape) != layer:
                continue
            if name is not None and name not in shape.name:
                continue
            if shape_id is not None and _shape_id(shape) != shape_id:
                continue
            if position is not None and shape_position != position:
                continue
            found.append((slide, shape))
    if not found:
        raise PatchApplyError(-1, f"no shapes matched select {select!r}")
    return found


def _rgb(value: str):
    from pptx.dml.color import RGBColor

    hex_value = value.strip().lstrip("#")
    if len(hex_value) != 6 or not re.fullmatch(r"[0-9A-Fa-f]{6}", hex_value):
        raise PatchApplyError(-1, f"invalid hex color {value!r} (expected RRGGBB)")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def _slide_by_index(prs, index: int):
    if not isinstance(index, int) or not 1 <= index <= len(prs.slides._sldIdLst):
        raise PatchApplyError(-1, f"slide {index} out of range (deck has {len(prs.slides._sldIdLst)} slides)")
    return prs.slides[index - 1]


def _move_slide(prs, current: int, target: Any) -> None:
    xml_slides = prs.slides._sldIdLst
    total = len(xml_slides)
    if not 1 <= current <= total:
        raise PatchApplyError(-1, f"slide {current} out of range (deck has {total} slides)")
    element = list(xml_slides)[current - 1]
    xml_slides.remove(element)
    if target == "end":
        xml_slides.append(element)
    elif target == "start":
        xml_slides.insert(0, element)
    elif isinstance(target, int) and 1 <= target <= total:
        # Removing one element first, so target maps onto the remaining list.
        xml_slides.insert(target - 1, element)
    else:
        raise PatchApplyError(-1, f"invalid move target {target!r} (1-based int, 'start' or 'end')")


def _apply_operation(prs, index: int, op: dict[str, Any], deck_title: str) -> dict[str, Any]:
    name = op["op"]
    try:
        if name == "set_text":
            affected = 0
            for _slide, shape in _select(prs, op["select"]):
                if getattr(shape, "has_text_frame", False):
                    shape.text_frame.text = str(op["text"])
                    affected += 1
            if affected == 0:
                raise PatchApplyError(index, "no selected shape has a text frame")
            return {"index": index, "op": name, "ok": True, "affected": affected}

        if name == "set_notes":
            slide = _slide_by_index(prs, op["slide"])
            slide.notes_slide.notes_text_frame.text = str(op["text"])
            return {"index": index, "op": name, "ok": True, "affected": 1}

        if name == "set_image":
            affected = 0
            for slide, shape in _select(prs, op["select"]):
                if shape.shape_type not in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE):
                    continue
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                old_name = shape.name
                c_nv_pr = shape._element.nvPicPr.find(qn("p:cNvPr"))
                old_alt = c_nv_pr.get("descr", "") if c_nv_pr is not None else ""
                alt = op.get("alt", old_alt)
                picture = slide.shapes.add_picture(str(op["resource"]), left, top, width=width, height=height)
                picture.name = old_name
                picture.shadow.inherit = False
                if alt:
                    new_c_nv_pr = picture._element.nvPicPr.find(qn("p:cNvPr"))
                    new_c_nv_pr.set("descr", alt)
                shape._element.getparent().remove(shape._element)
                affected += 1
            if affected == 0:
                raise PatchApplyError(index, "no selected shape is a picture")
            return {"index": index, "op": name, "ok": True, "affected": affected}

        if name == "set_shape_fill":
            rgb = _rgb(op["color"])
            affected = 0
            for _slide, shape in _select(prs, op["select"]):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    continue
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb
                affected += 1
            if affected == 0:
                raise PatchApplyError(index, "no applicable shape selected")
            return {"index": index, "op": name, "ok": True, "affected": affected}

        if name == "set_shape_line":
            rgb = _rgb(op["color"])
            affected = 0
            for _slide, shape in _select(prs, op["select"]):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    continue
                shape.line.fill.solid()
                shape.line.color.rgb = rgb
                shape.line.width = Pt(1.5)
                affected += 1
            if affected == 0:
                raise PatchApplyError(index, "no applicable shape selected")
            return {"index": index, "op": name, "ok": True, "affected": affected}

        if name == "set_shape_visible":
            visible = bool(op["visible"])
            affected = 0
            for _slide, shape in _select(prs, op["select"]):
                element = shape._element
                if visible:
                    if element.get("hidden") is not None:
                        del element.attrib["hidden"]
                        affected += 1
                elif element.get("hidden") != "1":
                    element.set("hidden", "1")
                    affected += 1
            return {"index": index, "op": name, "ok": True, "affected": affected}

        if name == "delete_shape":
            affected = 0
            for _slide, shape in _select(prs, op["select"]):
                shape._element.getparent().remove(shape._element)
                affected += 1
            return {"index": index, "op": name, "ok": True, "affected": affected}

        if name == "add_slide":
            slide_ir = parse_slide(op["slide_ir"], 0)
            envelope = DeckEnvelope(
                schema_version=IR_VERSION,
                kind=DOCUMENT_KIND,
                document_id="ppt-patch",
                metadata=Metadata(title=deck_title),
                document=PresentationIr(slides=()),
            )
            compiler = DeckCompiler(envelope, prs=prs)
            total = len(prs.slides._sldIdLst) + 1
            compiler.compile_slide(slide_ir, index=total, total=total)
            position = op.get("position", "end")
            if position != "end":
                _move_slide(prs, total, position)
            return {"index": index, "op": name, "ok": True, "affected": 1}

        if name == "delete_slide":
            _slide_by_index(prs, op["slide"])
            xml_slides = prs.slides._sldIdLst
            xml_slides.remove(list(xml_slides)[op["slide"] - 1])
            return {"index": index, "op": name, "ok": True, "affected": 1}

        if name == "move_slide":
            _move_slide(prs, op["slide"], op["to"])
            return {"index": index, "op": name, "ok": True, "affected": 1}

        if name == "assign_layer":
            select = op["select"]
            result = assign_layer(
                prs,
                op["layer"],
                slides=_normalize_slides(prs, select.get("slide", select.get("slides"))),
                shape=select.get("shape"),
                shape_id=select.get("id"),
                all_shapes=bool(select.get("all_shapes", False)),
            )
            return {"index": index, "op": name, "ok": True, "affected": result["tagged"], "detail": result["shapes"]}

        raise PatchApplyError(index, f"unhandled operation {name!r}")
    except PatchApplyError as exc:
        if exc.operation_index < 0:
            # -1 is the sentinel used by selector/slide validation helpers;
            # the patch contract reports the real failing operation index.
            raise PatchApplyError(index, exc.message) from exc
        raise
    except Exception as exc:  # wrap as a structured per-op error
        raise PatchApplyError(index, str(exc)) from exc


def apply_patch(input_path: str | Path, patch: dict[str, Any], output_path: str | Path | None = None) -> dict[str, Any]:
    """Apply a validated ppt-patch document; the deck is saved only on success."""
    patch = parse_patch(patch)
    source = Path(input_path)
    if not source.is_file():
        raise PatchApplyError(-1, f"input not found: {source}")
    target = Path(output_path) if output_path else source.with_name(f"{source.stem}-patched{source.suffix}")
    prs = Presentation(str(source))
    deck_title = str((patch.get("metadata") or {}).get("title", ""))
    reports = [_apply_operation(prs, index, op, deck_title) for index, op in enumerate(patch["operations"])]
    target.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(target))
    return {
        "success": True,
        "input": str(source),
        "output": str(target),
        "operations": reports,
        "slides_after": len(prs.slides._sldIdLst),
    }
